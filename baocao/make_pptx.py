# -*- coding: utf-8 -*-
"""
Sinh bài thuyết trình PowerPoint từ báo cáo BC-KIỂM THỬ.docx
Đề tài: Kiểm thử một số chức năng của trang web bằng công cụ Playwright
Trọng tâm: kiểm thử & các lỗi phát hiện (mỗi lỗi 1 slide).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

# ----------------------------------------------------------------------------
# THEME
# ----------------------------------------------------------------------------
PRIMARY = RGBColor(0x14, 0x3D, 0x59)   # xanh than chủ đạo
ACCENT  = RGBColor(0x2E, 0xAD, 0x33)   # xanh lá Playwright
ACCENT2 = RGBColor(0xF4, 0x7A, 0x1F)   # cam nhấn
LIGHT   = RGBColor(0xF3, 0xF6, 0xF8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x22, 0x2B, 0x33)
GRAY    = RGBColor(0x5B, 0x66, 0x70)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREEN_OK= RGBColor(0x1E, 0x8E, 0x3E)

FONT = "Segoe UI"
FONT_H = "Segoe UI Semibold"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_page = 0

# ----------------------------------------------------------------------------
# HELPERS
# ----------------------------------------------------------------------------
def _set_run(r, text, size, color, bold=False, italic=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font

def _bold_runs(p, text, size, color, font=FONT, base_bold=False):
    for tok in re.split(r"(\*\*.+?\*\*)", text):
        if tok == "":
            continue
        r = p.add_run()
        if tok.startswith("**") and tok.endswith("**"):
            _set_run(r, tok[2:-2], size, color, bold=True, font=font)
        else:
            _set_run(r, tok, size, color, bold=base_bold, font=font)

def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tf

def add_footer(slide):
    global _page
    _page += 1
    rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), PRIMARY)
    tf = textbox(slide, Inches(0.4), SH - Inches(0.33), Inches(9), Inches(0.3), MSO_ANCHOR.MIDDLE)
    _set_run(tf.paragraphs[0].add_run(), "Kiểm thử web bằng Playwright", 9, WHITE)
    tf2 = textbox(slide, SW - Inches(1.3), SH - Inches(0.33), Inches(0.9), Inches(0.3), MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    _set_run(p2.add_run(), str(_page), 9, WHITE, bold=True)

def header(slide, kicker, title):
    rect(slide, 0, 0, Inches(0.22), SH, ACCENT)
    tf = textbox(slide, Inches(0.55), Inches(0.30), Inches(12.3), Inches(0.35))
    _set_run(tf.paragraphs[0].add_run(), kicker.upper(), 12, ACCENT, bold=True)
    tf2 = textbox(slide, Inches(0.5), Inches(0.62), Inches(12.4), Inches(0.9))
    _bold_runs(tf2.paragraphs[0], title, 27, PRIMARY, font=FONT_H, base_bold=True)
    rect(slide, Inches(0.55), Inches(1.55), Inches(2.0), Pt(3), ACCENT)

def bullets(slide, items, x, y, w, h, size=16, gap=8):
    tf = textbox(slide, x, y, w, h)
    first = True
    for it in items:
        text, level = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.line_spacing = 1.05
        if level == 0:
            _set_run(p.add_run(), "▸  ", size, ACCENT, bold=True)
            _bold_runs(p, text, size, DARK)
        else:
            _set_run(p.add_run(), "      •  ", size-2, GRAY)
            _bold_runs(p, text, size-2, GRAY)
    return tf

def card(slide, x, y, w, h, title, body, accent=ACCENT, icon="", tsize=14.5, bsize=11.5):
    rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD,0xE3,0xE8))
    rect(slide, x, y, w, Inches(0.07), accent)
    tf = textbox(slide, x + Inches(0.18), y + Inches(0.2), w - Inches(0.36), h - Inches(0.34))
    p = tf.paragraphs[0]
    if icon:
        _set_run(p.add_run(), icon + "  ", tsize+1, accent, bold=True)
    _bold_runs(p, title, tsize, PRIMARY, font=FONT_H, base_bold=True)
    for b in body:
        pp = tf.add_paragraph(); pp.space_before = Pt(3); pp.line_spacing = 1.0
        _bold_runs(pp, b, bsize, GRAY)

def image_placeholder(slide, x, y, w, h, caption):
    box = rect(slide, x, y, w, h, RGBColor(0xEF,0xF4,0xF1), line=ACCENT)
    tf = textbox(slide, x + Inches(0.2), y, w - Inches(0.4), h, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), "🖼  KHUNG CHÈN ẢNH", 14, ACCENT, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    _set_run(p2.add_run(), caption, 11.5, GRAY, italic=True)

def make_table(slide, data, x, y, w, h, header_fill=PRIMARY, font_size=11,
               col_widths=None, header_size=11.5):
    rows = len(data); cols = len(data[0])
    gt = slide.shapes.add_table(rows, cols, x, y, w, h).table
    try:
        gt.first_row = True; gt.horz_banding = True
    except Exception:
        pass
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            gt.columns[j].width = Emu(int(w * cw / total))
    for i, rowdata in enumerate(data):
        for j, val in enumerate(rowdata):
            cell = gt.cell(i, j)
            cell.margin_left = Pt(5); cell.margin_right = Pt(5)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                _set_run(p.add_run(), str(val), header_size, WHITE, bold=True)
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                col = DARK; txt = str(val); bold = False
                if txt == "Pass": col = GREEN_OK; bold = True
                elif txt == "Fail": col = RED; bold = True
                elif txt == "High": col = RED; bold = True
                elif txt == "Medium": col = ACCENT2; bold = True
                elif txt == "Low": col = GRAY
                _set_run(p.add_run(), txt, font_size, col, bold=bold)
                if j == 0 or txt in ("Pass","Fail","High","Medium","Low"):
                    p.alignment = PP_ALIGN.CENTER
    return gt

def sev_color(sev):
    return {"High": RED, "Medium": ACCENT2, "Low": GRAY}[sev]

def new_slide():
    return prs.slides.add_slide(BLANK)

# ============================================================================
# 1 — TITLE
# ============================================================================
s = new_slide()
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, 0, 0, SW, Inches(0.18), ACCENT)
rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), ACCENT)
rect(s, SW - Inches(3.2), -Inches(1.5), Inches(4.5), Inches(4.5), RGBColor(0x1B,0x4C,0x6E), shape=MSO_SHAPE.OVAL)
rect(s, SW - Inches(1.6), SH - Inches(2.2), Inches(3.2), Inches(3.2), RGBColor(0x1B,0x4C,0x6E), shape=MSO_SHAPE.OVAL)
tf = textbox(s, Inches(0.9), Inches(0.7), Inches(10), Inches(0.9))
_set_run(tf.paragraphs[0].add_run(), "HỌC VIỆN KỸ THUẬT MẬT MÃ", 17, WHITE, bold=True)
p2 = tf.add_paragraph(); _set_run(p2.add_run(), "KHOA CÔNG NGHỆ THÔNG TIN", 14, RGBColor(0xBF,0xD9,0xE6))
rect(s, Inches(0.95), Inches(2.55), Inches(1.6), Pt(4), ACCENT)
tf = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.4))
_set_run(tf.paragraphs[0].add_run(), "BÁO CÁO KIỂM THỬ PHẦN MỀM NHÚNG", 18, ACCENT, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
_set_run(p2.add_run(), "KIỂM THỬ MỘT SỐ CHỨC NĂNG CỦA", 33, WHITE, bold=True, font=FONT_H)
p3 = tf.add_paragraph()
_set_run(p3.add_run(), "TRANG WEB BẰNG CÔNG CỤ PLAYWRIGHT", 33, WHITE, bold=True, font=FONT_H)
tf = textbox(s, Inches(0.9), Inches(5.45), Inches(11), Inches(1.4))
p = tf.paragraphs[0]
_set_run(p.add_run(), "Nhóm thực hiện:  ", 14, RGBColor(0xBF,0xD9,0xE6), bold=True)
_set_run(p.add_run(), "Ngô Hồng Phong · Đỗ Minh Thuần · Nguyễn Huy Hoàng", 14, WHITE)
p2 = tf.add_paragraph(); p2.space_before = Pt(4)
_set_run(p2.add_run(), "Giảng viên hướng dẫn:  ", 14, RGBColor(0xBF,0xD9,0xE6), bold=True)
_set_run(p2.add_run(), "Cô Thái Thị Thanh Vân", 14, WHITE)
p3 = tf.add_paragraph(); p3.space_before = Pt(8)
_set_run(p3.add_run(), "Hà Nội — 2026", 13, ACCENT, bold=True, italic=True)

# ============================================================================
# 2 — NỘI DUNG
# ============================================================================
s = new_slide()
header(s, "Tổng quan", "Nội dung trình bày")
items = [
    ("**Chương 1** — Tổng quan công cụ kiểm thử Playwright", ACCENT, "01"),
    ("**Chương 2** — Xây dựng kế hoạch kiểm thử", PRIMARY, "02"),
    ("**Chương 3** — Thực nghiệm & Kết quả kiểm thử", ACCENT2, "03"),
    ("**Trọng tâm** — Phân tích 8 lỗi phát hiện được", RGBColor(0x6A,0x4C,0x93), "★"),
]
y = Inches(1.95)
for text, col, num in items:
    rect(s, Inches(0.9), y, Inches(0.95), Inches(0.95), col)
    tf = textbox(s, Inches(0.9), y, Inches(0.95), Inches(0.95), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), num, 26, WHITE, bold=True, font=FONT_H)
    tf = textbox(s, Inches(2.1), y, Inches(10.5), Inches(0.95), MSO_ANCHOR.MIDDLE)
    _bold_runs(tf.paragraphs[0], text, 19, DARK)
    y += Inches(1.18)
add_footer(s)

# ============================================================================
# 3 — ĐẶT VẤN ĐỀ
# ============================================================================
s = new_slide()
header(s, "Lời mở đầu", "Đặt vấn đề & Mục tiêu")
bullets(s, [
    "Trang web ngày càng phức tạp, đòi hỏi **chất lượng cao** — kiểm thử thủ công tốn nhiều thời gian, công sức và dễ sót lỗi.",
    "**Kiểm thử tự động** trở thành thành phần không thể thiếu trong quy trình phát triển và bảo trì phần mềm.",
    "Nhóm lựa chọn **Playwright** (Microsoft) — framework kiểm thử thế hệ mới — để tự động hóa kịch bản kiểm thử.",
    "**Mục tiêu:** đánh giá hoạt động của website, **phát hiện lỗi** và đề xuất giải pháp nâng cao chất lượng.",
], Inches(0.9), Inches(1.95), Inches(7.2), Inches(4.6), size=17, gap=16)
card(s, Inches(8.45), Inches(2.0), Inches(4.0), Inches(4.2),
     "Phạm vi báo cáo", [
         "**C1.** Tổng quan & quy trình Playwright",
         "**C2.** Lập kế hoạch, chiến lược & phân công kiểm thử",
         "**C3.** Thực nghiệm 64 ca kiểm thử trên website thực tế",
         "**★** Phân tích chi tiết 8 lỗi tìm được",
     ], accent=PRIMARY, icon="📌")
add_footer(s)

# ============================================================================
# CHƯƠNG 1 (5 slide)
# ============================================================================
# 4 — Playwright là gì
s = new_slide()
header(s, "Chương 1", "Playwright là gì?")
bullets(s, [
    "Framework kiểm thử tự động **mã nguồn mở** do **Microsoft** phát triển.",
    "Kiểm thử web hiện đại trên 3 nhân trình duyệt: **Chromium** (Chrome, Edge), **Gecko** (Firefox), **WebKit** (Safari).",
    "Thực thi kịch bản **song song**, đa nền tảng → tiết kiệm thời gian.",
    "Cơ chế **Auto-wait** thông minh giúp loại bỏ lỗi bất đồng bộ (Flaky Test).",
    "Hỗ trợ đa ngôn ngữ: **JS, TS, Python, Java, C#**.",
], Inches(0.9), Inches(1.95), Inches(6.5), Inches(4.5), size=16.5, gap=14)
image_placeholder(s, Inches(7.8), Inches(1.95), Inches(4.7), Inches(4.3),
    "Logo / sơ đồ tổng quan Playwright\n(chụp từ trang playwright.dev)")
add_footer(s)

# 5 — Thành phần chính
s = new_slide()
header(s, "Chương 1", "Các thành phần chính")
w = Inches(2.92); h = Inches(2.6); gx = Inches(0.18); x0 = Inches(0.9); y0 = Inches(2.0)
comps = [
    ("Codegen", "🎬", "Ghi lại & phát lại thao tác trên trình duyệt, tự sinh mã kịch bản.", ACCENT),
    ("BrowserContext", "🧩", "Phiên trình duyệt độc lập — chạy song song nhiều tài khoản, không xung đột.", PRIMARY),
    ("Test Runner", "⚙️", "Bộ thực thi: tương tác trực tiếp trình duyệt, đa tab, chặn & kiểm tra request.", ACCENT2),
    ("Trace Viewer", "🔍", "Gỡ lỗi: xem lại video, ảnh chụp, console, network khi test thất bại.", RGBColor(0x6A,0x4C,0x93)),
]
for i, (t, ic, b, col) in enumerate(comps):
    card(s, x0 + i*(w+gx), y0, w, h, t, [b], accent=col, icon=ic, bsize=12.5)
add_footer(s)

# 6 — Ưu & nhược điểm
s = new_slide()
header(s, "Chương 1", "Ưu điểm & Nhược điểm")
box = rect(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(4.4), WHITE, line=GREEN_OK)
rect(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(0.55), GREEN_OK)
tf = textbox(s, Inches(1.1), Inches(1.97), Inches(5.2), Inches(0.5), MSO_ANCHOR.MIDDLE)
_set_run(tf.paragraphs[0].add_run(), "✓  Ưu điểm", 16, WHITE, bold=True)
bullets(s, [
    "Tốc độ thực thi **vượt trội**, tiết kiệm thời gian & chi phí.",
    "**Auto-wait** chính xác, đáng tin cậy, giảm sót lỗi.",
    "Tự động hóa tác vụ **phức tạp**: API, đa tab, đa trình duyệt.",
    "Tích hợp **CI/CD**, tăng năng suất đội ngũ.",
    "Đa ngôn ngữ, **mở rộng** linh hoạt theo nghiệp vụ.",
], Inches(1.1), Inches(2.6), Inches(5.2), Inches(3.6), size=14, gap=10)
box = rect(s, Inches(6.85), Inches(1.95), Inches(5.6), Inches(4.4), WHITE, line=RED)
rect(s, Inches(6.85), Inches(1.95), Inches(5.6), Inches(0.55), RED)
tf = textbox(s, Inches(7.05), Inches(1.97), Inches(5.2), Inches(0.5), MSO_ANCHOR.MIDDLE)
_set_run(tf.paragraphs[0].add_run(), "✗  Nhược điểm", 16, WHITE, bold=True)
bullets(s, [
    "Yêu cầu **kỹ năng lập trình** cao — rào cản với Manual Tester.",
    "Chi phí **xây dựng & bảo trì** kịch bản ban đầu lớn.",
    "Không hỗ trợ trình duyệt cũ (IE); mobile chỉ **giả lập**.",
    "Chạy song song **tốn CPU/RAM**.",
    "Khó với **CAPTCHA**, sinh trắc học, hệ thống chống bot.",
], Inches(7.05), Inches(2.6), Inches(5.2), Inches(3.6), size=14, gap=10)
add_footer(s)

# 7 — Quy trình làm việc
s = new_slide()
header(s, "Chương 1", "Quy trình làm việc với Playwright")
steps = [
    ("1", "Xác định kịch bản kiểm thử", "Phân tích nghiệp vụ, xác định input & kết quả kỳ vọng."),
    ("2", "Khởi tạo & cấu hình", "Thiết lập trình duyệt, headless, timeout, retries."),
    ("3", "Phát triển mã kịch bản", "Codegen + viết tay, tổ chức theo Page Object Model."),
    ("4", "Thực thi & chạy song song", "Phân chia đa luồng CPU & đa trình duyệt."),
    ("5", "Phân tích & gỡ lỗi", "Trace Viewer: timeline, video, ảnh, console, network."),
    ("6", "Xuất báo cáo & CI/CD", "Báo cáo HTML; tích hợp GitHub Actions, Jenkins."),
]
w = Inches(3.78); h = Inches(2.0); gx = Inches(0.2); gy = Inches(0.28)
for i,(n,t,b) in enumerate(steps):
    r,c = divmod(i,3)
    x = Inches(0.9) + c*(w+gx); y = Inches(1.95) + r*(h+gy)
    rect(s, x, y, w, h, WHITE, line=RGBColor(0xDD,0xE3,0xE8))
    rect(s, x, y, Inches(0.6), h, ACCENT)
    tf = textbox(s, x, y, Inches(0.6), h, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), n, 24, WHITE, bold=True, font=FONT_H)
    tf = textbox(s, x+Inches(0.75), y+Inches(0.18), w-Inches(0.9), h-Inches(0.3))
    _bold_runs(tf.paragraphs[0], "**"+t+"**", 13.5, PRIMARY, font=FONT_H)
    pp = tf.add_paragraph(); pp.space_before=Pt(4); _set_run(pp.add_run(), b, 11, GRAY)
add_footer(s)

# 8 — So sánh công cụ
s = new_slide()
header(s, "Chương 1", "So sánh các công cụ kiểm thử tự động")
data = [
    ["Tiêu chí", "Playwright", "Selenium", "Cypress", "Katalon"],
    ["Ngôn ngữ", "JS, TS, Python, Java, C#", "Java, C#, Python, Ruby, JS", "JS, TS", "Groovy, Java"],
    ["Trình duyệt", "Chromium, WebKit, Firefox", "Gần như tất cả (cả IE)", "Nhân Chromium & Firefox", "Chrome, FF, Edge, Safari, IE"],
    ["Kiến trúc", "WebSocket/CDP — cực nhanh", "WebDriver API trung gian", "Chạy trong Event Loop", "Dựa trên Selenium + Appium"],
    ["Cài đặt", "Dễ — 1 lệnh NPM", "Trung bình — quản lý Driver", "Dễ — qua NPM", "Dễ — đóng gói sẵn"],
    ["Chi phí", "Miễn phí, mã nguồn mở", "Miễn phí, mã nguồn mở", "Miễn phí (bản core)", "Nâng cao phải trả phí"],
]
make_table(s, data, Inches(0.9), Inches(1.85), Inches(11.55), Inches(3.7),
           font_size=10.5, header_size=11.5, col_widths=[1.1,1.6,1.6,1.4,1.5])
tf = textbox(s, Inches(0.9), Inches(5.75), Inches(11.55), Inches(1.3))
p = tf.paragraphs[0]
_set_run(p.add_run(), "→ Kết luận:  ", 14, ACCENT, bold=True)
_bold_runs(p, "Playwright là giải pháp **tối ưu** cho web hiện đại — tốc độ kiến trúc mới, đa trình duyệt mạnh mẽ, miễn phí và cài đặt tinh gọn.", 14, DARK)
add_footer(s)

# ============================================================================
# CHƯƠNG 2 (4 slide)
# ============================================================================
# 9 — Kế hoạch & website
s = new_slide()
header(s, "Chương 2", "Kế hoạch kiểm thử & Website thử nghiệm")
bullets(s, [
    "**Mục đích:** đảm bảo tính ổn định, kiểm thử **tương thích chéo** (cross-browser) và tối ưu năng suất.",
    "**Website thử nghiệm:** nền tảng khóa học trực tuyến",
    ("khoahoc.phongngohong.online", 1),
    "Đặc điểm: tải dữ liệu **động bất đồng bộ**, SPA, cây DOM phức tạp → dễ gây Flaky Test.",
    "→ Phù hợp để chứng minh sức mạnh **Auto-wait** của Playwright.",
], Inches(0.9), Inches(1.95), Inches(6.6), Inches(4.5), size=16, gap=14)
image_placeholder(s, Inches(7.85), Inches(1.95), Inches(4.6), Inches(4.3),
    "Ảnh chụp giao diện trang chủ website thử nghiệm\n(danh sách khóa học, thanh tìm kiếm, bộ lọc)")
add_footer(s)

# 10 — Phạm vi kiểm thử
s = new_slide()
header(s, "Chương 2", "Phạm vi kiểm thử")
rect(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(4.4), WHITE, line=GREEN_OK)
rect(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(0.55), GREEN_OK)
tf = textbox(s, Inches(1.1), Inches(1.97), Inches(5.2), Inches(0.5), MSO_ANCHOR.MIDDLE)
_set_run(tf.paragraphs[0].add_run(), "✓  Chức năng ĐƯỢC kiểm thử", 16, WHITE, bold=True)
bullets(s, [
    "**Tìm kiếm & Lọc khóa học** (từ khóa, danh mục, giá, phân trang).",
    "**Giỏ hàng & Đăng ký học** (thêm/xóa, mã giảm giá, tổng tiền realtime).",
    "**Đánh giá khóa học** (sao + bình luận).",
    "**Danh sách yêu thích** (Wishlist).",
], Inches(1.1), Inches(2.65), Inches(5.2), Inches(3.6), size=14, gap=12)
rect(s, Inches(6.85), Inches(1.95), Inches(5.6), Inches(4.4), WHITE, line=RED)
rect(s, Inches(6.85), Inches(1.95), Inches(5.6), Inches(0.55), RED)
tf = textbox(s, Inches(7.05), Inches(1.97), Inches(5.2), Inches(0.5), MSO_ANCHOR.MIDDLE)
_set_run(tf.paragraphs[0].add_run(), "✗  Chức năng KHÔNG kiểm thử", 16, WHITE, bold=True)
bullets(s, [
    "**Thanh toán tiền thật:** chỉ dừng ở nhập form, không kích hoạt giao dịch.",
    "**Xem video & tương tác:** streaming, tiến độ học, thi trắc nghiệm.",
    "**Quản trị hệ thống (Admin).**",
], Inches(7.05), Inches(2.65), Inches(5.2), Inches(3.6), size=14, gap=12)
add_footer(s)

# 11 — Lịch trình & nhân sự
s = new_slide()
header(s, "Chương 2", "Lịch trình công việc & Phân công nhân sự")
data = [
    ["Giai đoạn công việc", "Thời gian"],
    ["Nghiên cứu yêu cầu & Khảo sát Website", "11–13/05"],
    ["Xây dựng Kế hoạch kiểm thử (Test Plan)", "14–18/05"],
    ["Thiết kế Kịch bản kiểm thử chi tiết", "19–23/05"],
    ["Phát triển mã kịch bản Playwright", "24–28/05"],
    ["Chạy thực nghiệm, Gỡ lỗi & Tối ưu", "29/05–04/06"],
    ["Đánh giá kết quả & Tổng hợp báo cáo", "05–07/06"],
]
make_table(s, data, Inches(0.9), Inches(1.9), Inches(6.5), Inches(4.0),
           font_size=11.5, col_widths=[3.0, 1.0])
roles = [
    ("Ngô Hồng Phong", "Test Manager / Tester", "Quản lý tiến độ; kịch bản Giỏ hàng & Đăng ký."),
    ("Đỗ Minh Thuần", "Test Designer / Tester", "Thiết kế Test Cases; kịch bản Tìm kiếm."),
    ("Nguyễn Huy Hoàng", "Quality Analyst / Debugger", "Tối ưu POM, chạy song song, gỡ lỗi Trace Viewer."),
]
y = Inches(1.9)
for name, role, task in roles:
    card(s, Inches(7.7), y, Inches(4.75), Inches(1.28), name + "  —  " + role,
         [task], accent=ACCENT, icon="👤", tsize=13, bsize=11.5)
    y += Inches(1.4)
add_footer(s)

# 12 — Chiến lược & điều kiện chấp nhận
s = new_slide()
header(s, "Chương 2", "Chiến lược kiểm thử & Điều kiện chấp nhận")
card(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(4.35), "Chiến lược kiểm thử", [
    "**Kiểm thử hộp đen** (Black-box) tự động toàn phần.",
    "**Phân vùng tương đương** & **phân tích giá trị biên**.",
    "Dữ liệu đầu vào: hợp lệ, trống, ký tự đặc biệt, chuỗi quá dài, mã độc **SQLi / XSS**.",
    "Tự động hóa **100%** trên Playwright Test Runner.",
    "Chạy **Headless** (CI/CD) hoặc **Headed / UI Mode**.",
], accent=PRIMARY, icon="🎯", bsize=12.5)
card(s, Inches(6.85), Inches(1.95), Inches(5.6), Inches(4.35), "Điều kiện chấp nhận", [
    "Hiện thực hóa ≥ **95%** số Test Cases đã thiết kế.",
    "100% theo mô hình **POM**, dùng locator bền vững — không XPath tuyệt đối.",
    "Tỷ lệ **Pass ≥ 90%** trên cả 3 nhân trình duyệt.",
    "Loại bỏ hoàn toàn **Flaky Test** (≥ 3 phiên ổn định).",
    "Xuất đầy đủ **Trace log .zip** + báo cáo **HTML**.",
], accent=ACCENT, icon="✅", bsize=12.5)
add_footer(s)

# ============================================================================
# CHƯƠNG 3 (2 slide)
# ============================================================================
# 13 — Cách viết & thực thi kịch bản
s = new_slide()
header(s, "Chương 3", "Cách xây dựng & thực thi một kịch bản")
bullets(s, [
    "Mỗi ca kiểm thử theo cấu trúc **Arrange – Act – Assert**:",
    ("Reset dữ liệu về trạng thái ban đầu (before each).", 1),
    ("Thực hiện thao tác: mở trang, nhập liệu, click.", 1),
    ("Khẳng định (expect) kết quả mong đợi.", 1),
    "Thực thi: `npx playwright test` chạy lần lượt **64 test case**.",
    "Test **PASS** ✓ (xanh) / **FAIL** ✘ (đỏ) kèm lỗi; tự **chụp ảnh** khi fail.",
], Inches(0.9), Inches(1.95), Inches(6.4), Inches(4.4), size=15, gap=11)
image_placeholder(s, Inches(7.65), Inches(1.95), Inches(4.8), Inches(4.3),
    "Ảnh chụp kết quả chạy test trên terminal:\ncác test PASS (xanh ✓) và FAIL (đỏ ✘)")
add_footer(s)

# 14 — Kết quả tổng hợp
s = new_slide()
header(s, "Chương 3", "Kết quả kiểm thử — 64 Test Case")
data = [
    ["Chức năng", "Pass", "Fail", "Tổng", "Tỉ lệ"],
    ["Tìm kiếm & Lọc khóa học", "8", "2", "10", "80.00%"],
    ["Đánh giá khóa học", "13", "3", "16", "81.25%"],
    ["Đăng ký & Thanh toán", "20", "2", "22", "90.91%"],
    ["Danh sách yêu thích", "15", "1", "16", "93.75%"],
    ["TỔNG CỘNG", "56", "8", "64", "87.50%"],
]
make_table(s, data, Inches(0.9), Inches(2.0), Inches(6.8), Inches(3.4),
           font_size=12.5, header_size=13, col_widths=[2.6,0.7,0.7,0.8,1.0])
kpis = [("64","Tổng test case",PRIMARY),("56","Test PASS",GREEN_OK),("8","Test FAIL",RED),("87.5%","Tỉ lệ đạt",ACCENT)]
w = Inches(2.25)
for i,(num,lab,col) in enumerate(kpis):
    r,c = divmod(i,2)
    x = Inches(8.0) + c*(w+Inches(0.3)); y = Inches(2.0) + r*(Inches(1.55)+Inches(0.25))
    rect(s, x, y, w, Inches(1.55), col)
    tf = textbox(s, x, y+Inches(0.15), w, Inches(1.3), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    _set_run(p.add_run(), num, 32, WHITE, bold=True, font=FONT_H)
    pp = tf.add_paragraph(); pp.alignment=PP_ALIGN.CENTER
    _set_run(pp.add_run(), lab, 13, WHITE)
tf = textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.9))
_bold_runs(tf.paragraphs[0], "Thiết kế & thực thi thành công **64 ca kiểm thử** cho **4 nhóm chức năng** — phát hiện **8 lỗi** thực tế.", 15, DARK)
add_footer(s)

# ============================================================================
# SECTION — 8 LỖI (mỗi lỗi 1 slide)
# ============================================================================
# 15 — Phân loại & tổng quan lỗi
s = new_slide()
header(s, "Phân tích lỗi", "Tổng quan 8 lỗi phát hiện được")
levels = [
    ("Cao (High) — 3 lỗi", "Tính năng cốt lõi tê liệt hoặc rủi ro bảo mật / hệ thống sập.", RED),
    ("Trung bình (Medium) — 3 lỗi", "Sai lệch logic nghiệp vụ nhưng không làm tê liệt hệ thống.", ACCENT2),
    ("Thấp (Low) — 2 lỗi", "Lỗi hiển thị, định dạng số — ảnh hưởng trải nghiệm.", GRAY),
]
w = Inches(3.78); gx = Inches(0.2)
for i,(t,b,col) in enumerate(levels):
    card(s, Inches(0.9)+i*(w+gx), Inches(1.95), w, Inches(1.9), t, [b], accent=col, icon="🐞", tsize=13)
data = [
    ["TC", "Lỗi", "Nhóm chức năng", "Mức"],
    ["TC-09", "Tràn độ dài từ khóa → HTTP 500", "Tìm kiếm", "High"],
    ["TC-24", "Lỗ hổng XSS trong bình luận", "Đánh giá", "High"],
    ["TC-59", "Spam nút yêu thích → HTTP 500", "Wishlist", "High"],
    ["TC-61", "Lọc Danh mục + Giá sai kết quả", "Tìm kiếm/Lọc", "Medium"],
    ["TC-23", "Race condition — đánh giá trùng", "Đánh giá", "Medium"],
    ["TC-64", "Ẩn khóa đang chờ thanh toán", "Đăng ký", "Medium"],
    ["TC-62", "Điểm trung bình không làm tròn", "Đánh giá", "Low"],
    ["TC-63", "Sai số làm tròn tiền giảm giá", "Đăng ký", "Low"],
]
make_table(s, data, Inches(0.9), Inches(4.05), Inches(11.55), Inches(2.2),
           font_size=11, header_size=11.5, col_widths=[0.8, 3.0, 1.4, 0.7])
add_footer(s)

# ----- dữ liệu chi tiết 8 lỗi -----
bugs = [
    dict(tc="TC-09", title="Tràn độ dài từ khóa tìm kiếm", area="Tìm kiếm khóa học", sev="High",
         tech="Phân tích giá trị biên (Boundary Value)",
         phenom="Nhập từ khóa rất dài (chuỗi 210 ký tự) vào ô tìm kiếm. Frontend không giới hạn độ dài đầu vào.",
         expected="Hệ thống báo lỗi thân thiện, từ chối truy vấn an toàn.",
         actual="Backend nhận chuỗi > 200 ký tự và trả về lỗi HTTP 500 (sập xử lý).",
         cause="Thiếu kiểm tra ràng buộc độ dài ở cả phía nhập liệu và xử lý truy vấn."),
    dict(tc="TC-24", title="Lỗ hổng XSS trong bình luận đánh giá", area="Đánh giá khóa học", sev="High",
         tech="Kiểm thử bảo mật — chèn mã độc (XSS)",
         phenom="Gửi bình luận chứa mã script độc hại, sau đó tải lại trang hiển thị đánh giá.",
         expected="Mã script bị vô hiệu hóa, chỉ hiển thị dưới dạng văn bản thuần.",
         actual="Mã script độc hại được trình duyệt thực thi → lỗ hổng XSS.",
         cause="Hiển thị nội dung qua dangerouslySetInnerHTML mà không lọc/escape HTML."),
    dict(tc="TC-59", title="Spam nút yêu thích gây sập API", area="Danh sách yêu thích", sev="High",
         tech="Kiểm thử tương tranh (Concurrency / Spam)",
         phenom="Bấm liên tục thật nhanh vào nút yêu thích của cùng một khóa học.",
         expected="Hệ thống xử lý ổn định, không phát sinh lỗi.",
         actual="Nhiều request đồng thời gây xung đột ràng buộc UNIQUE → trả về HTTP 500.",
         cause="Frontend không debounce nút; backend chưa xử lý xung đột khóa duy nhất."),
    dict(tc="TC-61", title="Lọc kết hợp Danh mục + Giá sai kết quả", area="Tìm kiếm / Lọc", sev="Medium",
         tech="Kiểm thử tổ hợp điều kiện bộ lọc",
         phenom="Chọn đồng thời danh mục \"Go\" và mức giá \"Có phí\".",
         expected="Hiển thị đúng 1 khóa (Go và có phí).",
         actual="Trả về 2 khóa — bộ lọc giá bị bỏ qua khi đã chọn danh mục.",
         cause="Logic backend dùng else-if giữa điều kiện danh mục và giá thay vì kết hợp cả hai."),
    dict(tc="TC-23", title="Race condition tạo đánh giá trùng", area="Đánh giá khóa học", sev="Medium",
         tech="Kiểm thử tương tranh (Race Condition)",
         phenom="Bấm nút \"Gửi đánh giá\" nhiều lần liên tiếp trong thời gian ngắn.",
         expected="Chỉ tồn tại duy nhất 1 đánh giá của người dùng.",
         actual="Tạo ra 2 đánh giá trùng của cùng một người.",
         cause="Frontend không disable nút Gửi; backend ghi dữ liệu có độ trễ."),
    dict(tc="TC-64", title="Ẩn khóa học đang chờ thanh toán", area="Đăng ký & Thanh toán", sev="Medium",
         tech="Kiểm thử trạng thái dữ liệu",
         phenom="Đăng ký khóa trả phí (trạng thái chờ thanh toán) rồi mở \"Khóa học của tôi\".",
         expected="Khóa học hiển thị trong danh sách của người dùng.",
         actual="Khóa bị ẩn hoàn toàn khỏi \"Khóa học của tôi\".",
         cause="API my-courses chỉ lọc status='completed', bỏ sót trạng thái đang chờ."),
    dict(tc="TC-62", title="Điểm trung bình không làm tròn", area="Đánh giá khóa học", sev="Low",
         tech="Kiểm thử hiển thị & định dạng số",
         phenom="Xem điểm đánh giá trung bình tại trang chi tiết khóa học.",
         expected="Hiển thị làm tròn 1 chữ số thập phân (vd: 3.3).",
         actual="Hiển thị 3.3333333333333335 (số thập phân thô).",
         cause="Trang chi tiết quên gọi hàm làm tròn .toFixed(1)."),
    dict(tc="TC-63", title="Sai số làm tròn tiền giảm giá", area="Đăng ký & Thanh toán", sev="Low",
         tech="Kiểm thử tính toán & làm tròn số tiền",
         phenom="Áp mã giảm giá 20% cho khóa học giá 49.99.",
         expected="Hiển thị số tiền làm tròn 2 chữ số (39.99).",
         actual="Hiển thị 39.992000000000004 — sai số dấu phẩy động.",
         cause="Backend tính 49.99 * 0.8 bằng float, không làm tròn kết quả."),
]

total_bugs = len(bugs)
for idx, b in enumerate(bugs, 1):
    s = new_slide()
    col = sev_color(b["sev"])
    header(s, f"Phân tích lỗi  ·  {idx}/{total_bugs}", f"{b['tc']} — {b['title']}")
    # badge mức độ + chức năng
    badge = rect(s, Inches(9.95), Inches(0.45), Inches(2.5), Inches(0.7), col)
    tf = textbox(s, Inches(9.95), Inches(0.45), Inches(2.5), Inches(0.7), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), "MỨC ĐỘ: " + b["sev"].upper(), 13, WHITE, bold=True)
    # chức năng + kỹ thuật phát hiện
    tf = textbox(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.5))
    p = tf.paragraphs[0]
    _set_run(p.add_run(), "Nhóm chức năng:  ", 13, ACCENT, bold=True)
    _set_run(p.add_run(), b["area"] + "      ", 13, DARK)
    _set_run(p.add_run(), "Kỹ thuật phát hiện:  ", 13, ACCENT, bold=True)
    _set_run(p.add_run(), b["tech"], 13, DARK)
    # hiện tượng
    card(s, Inches(0.9), Inches(2.35), Inches(11.55), Inches(1.25), "Hiện tượng / Cách kiểm thử",
         [b["phenom"]], accent=PRIMARY, icon="🔬", tsize=14, bsize=13)
    # mong đợi vs thực tế
    rect(s, Inches(0.9), Inches(3.8), Inches(5.72), Inches(1.55), WHITE, line=GREEN_OK)
    rect(s, Inches(0.9), Inches(3.8), Inches(5.72), Inches(0.45), GREEN_OK)
    tf = textbox(s, Inches(1.05), Inches(3.8), Inches(5.5), Inches(0.45), MSO_ANCHOR.MIDDLE)
    _set_run(tf.paragraphs[0].add_run(), "✓  Kết quả mong đợi", 13, WHITE, bold=True)
    tf = textbox(s, Inches(1.05), Inches(4.3), Inches(5.4), Inches(1.0))
    _bold_runs(tf.paragraphs[0], b["expected"], 12.5, DARK)
    rect(s, Inches(6.73), Inches(3.8), Inches(5.72), Inches(1.55), WHITE, line=RED)
    rect(s, Inches(6.73), Inches(3.8), Inches(5.72), Inches(0.45), RED)
    tf = textbox(s, Inches(6.88), Inches(3.8), Inches(5.5), Inches(0.45), MSO_ANCHOR.MIDDLE)
    _set_run(tf.paragraphs[0].add_run(), "✗  Kết quả thực tế (LỖI)", 13, WHITE, bold=True)
    tf = textbox(s, Inches(6.88), Inches(4.3), Inches(5.4), Inches(1.0))
    _bold_runs(tf.paragraphs[0], b["actual"], 12.5, RED)
    # nguyên nhân
    rect(s, Inches(0.9), Inches(5.55), Inches(11.55), Inches(0.95), LIGHT)
    rect(s, Inches(0.9), Inches(5.55), Inches(0.1), Inches(0.95), col)
    tf = textbox(s, Inches(1.2), Inches(5.55), Inches(11.0), Inches(0.95), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _set_run(p.add_run(), "Nguyên nhân:  ", 13, col, bold=True)
    _bold_runs(p, b["cause"], 13, DARK)
    image_placeholder(s, Inches(0.9), Inches(6.62), Inches(11.55), Inches(0.5),
        "(Tùy chọn) Dán ảnh chụp minh họa lỗi: terminal / màn hình / Trace Viewer của " + b["tc"])
    add_footer(s)

# ============================================================================
# HẠN CHẾ & HƯỚNG PHÁT TRIỂN
# ============================================================================
s = new_slide()
header(s, "Kết luận", "Hạn chế & Hướng phát triển")
card(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(4.35), "Hạn chế", [
    "Mới kiểm thử **4 nhóm chức năng** cốt lõi, chưa bao phủ phần **Admin**.",
    "Chưa kiểm thử **hiệu năng (performance)** và **tải (load testing)**.",
    "Chưa khai thác hết tính năng nâng cao: chạy **đa trình duyệt** song song, **visual testing**, tích hợp **CI/CD** đầy đủ.",
], accent=ACCENT2, icon="⚠️", bsize=13)
card(s, Inches(6.85), Inches(1.95), Inches(5.6), Inches(4.35), "Hướng phát triển", [
    "Mở rộng phạm vi kiểm thử cho **toàn bộ hệ thống** (gồm Admin).",
    "Tích hợp bộ test vào **CI/CD** (GitHub Actions) — tự chạy khi đổi mã nguồn.",
    "Kiểm thử **đa trình duyệt** (Firefox, WebKit) & **thiết bị di động**.",
    "Áp dụng **visual regression** và kiểm thử **hiệu năng**.",
], accent=ACCENT, icon="🚀", bsize=13)
add_footer(s)

# ============================================================================
# CẢM ƠN
# ============================================================================
s = new_slide()
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, 0, 0, SW, Inches(0.18), ACCENT)
rect(s, 0, SH-Inches(0.18), SW, Inches(0.18), ACCENT)
rect(s, SW-Inches(3.0), -Inches(1.4), Inches(4.2), Inches(4.2), RGBColor(0x1B,0x4C,0x6E), shape=MSO_SHAPE.OVAL)
tf = textbox(s, Inches(0), Inches(2.5), SW, Inches(1.6), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "CẢM ƠN THẦY CÔ ĐÃ LẮNG NGHE!", 40, WHITE, bold=True, font=FONT_H)
tf = textbox(s, Inches(0), Inches(4.25), SW, Inches(1.2), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "Nhóm xin chân thành cảm ơn cô Thái Thị Thanh Vân", 17, RGBColor(0xCF,0xE3,0xEC))
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
_set_run(p2.add_run(), "Ngô Hồng Phong · Đỗ Minh Thuần · Nguyễn Huy Hoàng", 15, ACCENT, bold=True)

prs.save("BaoCao_Playwright_Slide.pptx")
print("SAVED -> BaoCao_Playwright_Slide.pptx | slides:", len(prs.slides._sldIdLst))
